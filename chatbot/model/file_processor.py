import os
import logging
from typing import Dict, Any, Optional, List
import chromadb
from google import genai
from google.genai import types
import PyPDF2
import docx
import openpyxl
import pptx
from io import BytesIO
import hashlib
from dotenv import load_dotenv

# Load environment variables
parent_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
dotenv_path = os.path.join(parent_dir, 'backend', '.env')
load_dotenv(dotenv_path=dotenv_path)

logging.basicConfig(level=logging.INFO)

class FileProcessor:
    def __init__(self):
        """Khởi tạo FileProcessor với Gemini và ChromaDB"""
        try:
            # Khởi tạo Gemini Client
            api_key = os.getenv("GEMINI_API_KEY")
            if not api_key:
                raise ValueError("GEMINI_API_KEY not set")

            self.client = genai.Client(api_key=api_key)
            self.model = "gemini-2.0-flash-exp"
            self.embedding_model = 'text-embedding-004'

            # Khởi tạo ChromaDB cho files
            self.chroma_path = "chroma_db_files"
            self.chroma_client = chromadb.PersistentClient(path=self.chroma_path)

            # Tạo collection cho files
            try:
                self.file_collection = self.chroma_client.get_collection(name="uploaded_files")
            except:
                self.file_collection = self.chroma_client.create_collection(
                    name="uploaded_files",
                    metadata={"hnsw:space": "cosine"}
                )

            logging.info("FileProcessor initialized successfully")

        except Exception as e:
            logging.error(f"Failed to initialize FileProcessor: {e}")
            raise

    def extract_text_from_pdf(self, file_content: BytesIO) -> str:
        """Trích xuất văn bản từ file PDF"""
        try:
            pdf_reader = PyPDF2.PdfReader(file_content)
            text = ""
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            logging.error(f"Error extracting PDF text: {e}")
            return ""

    def extract_text_from_docx(self, file_content: BytesIO) -> str:
        """Trích xuất văn bản từ file DOCX"""
        try:
            doc = docx.Document(file_content)
            text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text.append(cell.text)

            return "\n".join(text)
        except Exception as e:
            logging.error(f"Error extracting DOCX text: {e}")
            return ""

    def extract_text_from_xlsx(self, file_content: BytesIO) -> str:
        """Trích xuất văn bản từ file Excel"""
        try:
            workbook = openpyxl.load_workbook(file_content, read_only=True)
            text = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text.append(f"Sheet: {sheet_name}")

                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        text.append(row_text)

            return "\n".join(text)
        except Exception as e:
            logging.error(f"Error extracting XLSX text: {e}")
            return ""

    def extract_text_from_pptx(self, file_content: BytesIO) -> str:
        """Trích xuất văn bản từ file PowerPoint"""
        try:
            presentation = pptx.Presentation(file_content)
            text = []

            for slide_num, slide in enumerate(presentation.slides, 1):
                text.append(f"Slide {slide_num}:")
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        if shape.text.strip():
                            text.append(shape.text)

            return "\n".join(text)
        except Exception as e:
            logging.error(f"Error extracting PPTX text: {e}")
            return ""

    def extract_text_from_txt(self, file_content: BytesIO) -> str:
        """Trích xuất văn bản từ file text"""
        try:
            text = file_content.read()
            if isinstance(text, bytes):
                text = text.decode('utf-8', errors='ignore')
            return text
        except Exception as e:
            logging.error(f"Error extracting TXT text: {e}")
            return ""

    def extract_text(self, file_content: BytesIO, filename: str, content_type: str) -> str:
        """Trích xuất văn bản từ file dựa trên loại file"""
        file_content.seek(0)  # Reset position

        # Xác định loại file
        ext = os.path.splitext(filename)[1].lower()

        if ext == '.pdf' or content_type == 'application/pdf':
            return self.extract_text_from_pdf(file_content)
        elif ext in ['.docx', '.doc'] or 'word' in content_type:
            return self.extract_text_from_docx(file_content)
        elif ext in ['.xlsx', '.xls'] or 'excel' in content_type or 'spreadsheet' in content_type:
            return self.extract_text_from_xlsx(file_content)
        elif ext in ['.pptx', '.ppt'] or 'powerpoint' in content_type or 'presentation' in content_type:
            return self.extract_text_from_pptx(file_content)
        elif ext == '.txt' or 'text/plain' in content_type:
            return self.extract_text_from_txt(file_content)
        else:
            logging.warning(f"Unsupported file type: {filename}")
            return ""

    def chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Chia văn bản thành các chunks nhỏ hơn"""
        if not text:
            return []

        chunks = []
        start = 0
        text_length = len(text)

        while start < text_length:
            end = start + chunk_size

            # Tìm điểm ngắt tự nhiên (dấu chấm, xuống dòng)
            if end < text_length:
                for delimiter in ['\n\n', '\n', '. ', '! ', '? ']:
                    delimiter_pos = text.rfind(delimiter, start + overlap, end)
                    if delimiter_pos != -1:
                        end = delimiter_pos + len(delimiter)
                        break

            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)

            start = end - overlap if end < text_length else end

        return chunks

    def summarize_text(self, text: str, max_length: int = 500) -> str:
        """Tóm tắt văn bản sử dụng Gemini"""
        try:
            if not text:
                return "Không có nội dung để tóm tắt"

            # Giới hạn văn bản đầu vào để tránh quá tải
            max_input_length = 10000
            if len(text) > max_input_length:
                text = text[:max_input_length] + "..."

            prompt = f"""Hãy tóm tắt nội dung sau đây một cách ngắn gọn, súc tích trong khoảng {max_length} ký tự.
            Tập trung vào các ý chính và thông tin quan trọng nhất.
            
            Nội dung:
            {text}
            
            Tóm tắt:"""

            response = self.client.models.generate_content(
                model=self.model,
                contents=prompt,
                config=types.GenerateContentConfig(
                    temperature=0.3,
                    max_output_tokens=max_length // 4,  # Ước lượng tokens
                )
            )

            return response.text.strip()

        except Exception as e:
            logging.error(f"Error summarizing text: {e}")
            return "Lỗi khi tóm tắt nội dung"

    def create_embeddings(self, texts: List[str]) -> List[List[float]]:
        """Tạo embeddings cho danh sách văn bản"""
        embeddings = []

        for text in texts:
            try:
                response = self.client.models.embed_content(
                    model=self.embedding_model,
                    contents=[text]
                )
                embeddings.append(response.values[0])
            except Exception as e:
                logging.error(f"Error creating embedding: {e}")
                # Thêm vector rỗng nếu lỗi
                embeddings.append([0.0] * 768)  # Kích thước vector mặc định

        return embeddings

    def process_file(self, file_content: BytesIO, filename: str, content_type: str, file_id: str) -> Dict[str, Any]:
        """Xử lý file: trích xuất nội dung, tóm tắt và vector hóa"""
        try:
            # 1. Trích xuất văn bản
            text = self.extract_text(file_content, filename, content_type)

            if not text:
                return {
                    'success': False,
                    'error': 'Không thể trích xuất nội dung từ file'
                }

            # 2. Tóm tắt nội dung
            summary = self.summarize_text(text, max_length=500)

            # 3. Chia thành chunks
            chunks = self.chunk_text(text)

            if not chunks:
                chunks = [text[:1000]]  # Lấy 1000 ký tự đầu nếu không chia được

            # 4. Tạo embeddings
            embeddings = self.create_embeddings(chunks)

            # 5. Tạo ID duy nhất cho vector
            vector_id = f"file_{file_id}_{hashlib.md5(filename.encode()).hexdigest()[:8]}"

            # 6. Lưu vào ChromaDB
            ids = [f"{vector_id}_{i}" for i in range(len(chunks))]
            metadatas = [
                {
                    'file_id': file_id,
                    'filename': filename,
                    'chunk_index': i,
                    'total_chunks': len(chunks)
                }
                for i in range(len(chunks))
            ]

            self.file_collection.add(
                embeddings=embeddings,
                documents=chunks,
                metadatas=metadatas,
                ids=ids
            )

            logging.info(f"Successfully processed file: {filename} with {len(chunks)} chunks")

            return {
                'success': True,
                'content': text[:5000],  # Giới hạn nội dung trả về
                'summary': summary,
                'vector_id': vector_id,
                'chunks_count': len(chunks)
            }

        except Exception as e:
            logging.error(f"Error processing file: {e}")
            return {
                'success': False,
                'error': str(e)
            }

    def delete_vector(self, vector_id: str) -> bool:
        """Xóa vector khỏi ChromaDB"""
        try:
            # Lấy tất cả IDs liên quan đến vector_id
            results = self.file_collection.get(
                where={"file_id": vector_id.replace("file_", "")}
            )

            if results['ids']:
                self.file_collection.delete(ids=results['ids'])
                logging.info(f"Deleted {len(results['ids'])} vectors for {vector_id}")
                return True

            return True

        except Exception as e:
            logging.error(f"Error deleting vector: {e}")
            return False

    def search_in_files(self, query: str, n_results: int = 3) -> List[Dict[str, Any]]:
        """Tìm kiếm trong các file đã upload"""
        try:
            # Tạo embedding cho query
            query_embedding = self.client.models.embed_content(
                model=self.embedding_model,
                contents=query
            ).embedding

            # Tìm kiếm trong ChromaDB
            results = self.file_collection.query(
                query_embeddings=[query_embedding],
                n_results=n_results,
                include=['documents', 'metadatas', 'distances']
            )

            search_results = []
            if results['documents'] and len(results['documents'][0]) > 0:
                for doc, metadata, distance in zip(
                        results['documents'][0],
                        results['metadatas'][0],
                        results['distances'][0]
                ):
                    search_results.append({
                        'content': doc,
                        'filename': metadata.get('filename', 'Unknown'),
                        'file_id': metadata.get('file_id', ''),
                        'similarity': 1 - distance  # Chuyển distance thành similarity
                    })

            return search_results

        except Exception as e:
            logging.error(f"Error searching in files: {e}")
            return []

    def get_all_vectors_info(self) -> List[Dict[str, Any]]:
        """Lấy thông tin về tất cả vectors đã lưu"""
        try:
            # Lấy tất cả dữ liệu từ collection
            all_data = self.file_collection.get()

            # Group by file_id
            files_info = {}
            for metadata in all_data['metadatas']:
                file_id = metadata.get('file_id', '')
                if file_id not in files_info:
                    files_info[file_id] = {
                        'file_id': file_id,
                        'filename': metadata.get('filename', 'Unknown'),
                        'chunks_count': 0
                    }
                files_info[file_id]['chunks_count'] += 1

            return list(files_info.values())

        except Exception as e:
            logging.error(f"Error getting vectors info: {e}")
            return []